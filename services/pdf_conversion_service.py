import base64
import os
import subprocess
import tempfile
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors


ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt'}


def is_allowed_file(filename):
    ext = os.path.splitext(filename)[1].lower()
    return ext in ALLOWED_EXTENSIONS


def convert_to_pdf(src_path, dest_path):
    """Convert docx/xlsx/pptx to PDF. Returns (pdf_path, error)."""
    ext = os.path.splitext(src_path)[1].lower()
    if ext == '.pdf':
        if os.path.abspath(src_path) == os.path.abspath(dest_path):
            return dest_path, None
        import shutil
        shutil.copy2(src_path, dest_path)
        return dest_path, None
    try:
        if ext in ('.docx', '.doc'):
            return _docx_to_pdf(src_path, dest_path)
        elif ext in ('.xlsx', '.xls'):
            return _xlsx_to_pdf(src_path, dest_path)
        elif ext in ('.pptx', '.ppt'):
            return _pptx_to_pdf(src_path, dest_path)
        else:
            return None, f"Unsupported format: {ext}"
    except Exception as e:
        return None, str(e)


def _docx_to_pdf(src_path, dest_path):
    from docx import Document
    doc = Document(src_path)
    pdf_doc = SimpleDocTemplate(dest_path, pagesize=A4,
                                 leftMargin=inch, rightMargin=inch,
                                 topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    story = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            story.append(Spacer(1, 6))
            continue
        style = styles['Heading1'] if para.style.name.startswith('Heading 1') else \
                styles['Heading2'] if para.style.name.startswith('Heading 2') else \
                styles['Normal']
        try:
            story.append(Paragraph(text, style))
        except Exception:
            story.append(Paragraph(text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;'), style))
        story.append(Spacer(1, 4))
    if not story:
        story.append(Paragraph("(Empty document)", styles['Normal']))
    pdf_doc.build(story)
    return dest_path, None


def _xlsx_to_pdf(src_path, dest_path):
    from openpyxl import load_workbook
    wb = load_workbook(src_path, data_only=True)
    pdf_doc = SimpleDocTemplate(dest_path, pagesize=A4,
                                 leftMargin=0.5*inch, rightMargin=0.5*inch,
                                 topMargin=0.75*inch, bottomMargin=0.75*inch)
    styles = getSampleStyleSheet()
    story = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        story.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading2']))
        story.append(Spacer(1, 8))
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            story.append(Paragraph("(Empty sheet)", styles['Normal']))
            continue
        # Limit columns and rows for readability
        max_cols = min(len(rows[0]) if rows[0] else 0, 10)
        max_rows = min(len(rows), 200)
        table_data = []
        for row in rows[:max_rows]:
            table_data.append([str(cell) if cell is not None else '' for cell in row[:max_cols]])
        if table_data:
            col_width = (A4[0] - inch) / max(max_cols, 1)
            t = Table(table_data, colWidths=[col_width] * max_cols, repeatRows=1)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1c2128')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#f0a500')),
                ('FONTSIZE', (0, 0), (-1, -1), 7),
                ('GRID', (0, 0), (-1, -1), 0.3, colors.HexColor('#30363d')),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor('#161b22'), colors.HexColor('#0d1117')]),
                ('TEXTCOLOR', (0, 1), (-1, -1), colors.HexColor('#e6edf3')),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('WORDWRAP', (0, 0), (-1, -1), True),
            ]))
            story.append(t)
        story.append(Spacer(1, 16))
    pdf_doc.build(story)
    return dest_path, None


def _pptx_to_pdf(src_path, dest_path):
    """Render a presentation as real slides when Microsoft PowerPoint is available."""
    if os.name == 'nt':
        try:
            _pptx_to_pdf_with_powerpoint(src_path, dest_path)
            return dest_path, None
        except (OSError, RuntimeError, subprocess.SubprocessError):
            # Keep the text fallback for servers that do not have PowerPoint installed.
            pass

    return _pptx_to_text_pdf(src_path, dest_path)


def _pptx_to_pdf_with_powerpoint(src_path, dest_path):
    """Use PowerPoint's own PDF exporter so slide layouts and graphics are preserved."""
    source_path = os.path.abspath(src_path)
    destination_path = os.path.abspath(dest_path)
    destination_dir = os.path.dirname(destination_path) or os.getcwd()
    os.makedirs(destination_dir, exist_ok=True)

    temporary_file = tempfile.NamedTemporaryFile(
        prefix='dms-pptx-preview-',
        suffix='.pdf',
        dir=destination_dir,
        delete=False,
    )
    temporary_path = temporary_file.name
    temporary_file.close()
    os.remove(temporary_path)

    script = r'''
$ErrorActionPreference = 'Stop'
$sourcePath = [Environment]::GetEnvironmentVariable('DMS_PPTX_SOURCE')
$destinationPath = [Environment]::GetEnvironmentVariable('DMS_PPTX_DESTINATION')
$powerPoint = $null
$presentation = $null

try {
    $powerPoint = New-Object -ComObject PowerPoint.Application
    $presentation = $powerPoint.Presentations.Open($sourcePath, $true, $true, $false)
    $presentation.SaveAs($destinationPath, 32)
}
finally {
    if ($null -ne $presentation) {
        $presentation.Close()
        [System.Runtime.InteropServices.Marshal]::ReleaseComObject($presentation) | Out-Null
    }
    if ($null -ne $powerPoint) {
        $powerPoint.Quit()
        [System.Runtime.InteropServices.Marshal]::ReleaseComObject($powerPoint) | Out-Null
    }
    [GC]::Collect()
    [GC]::WaitForPendingFinalizers()
}
'''
    encoded_script = base64.b64encode(script.encode('utf-16-le')).decode('ascii')
    environment = os.environ.copy()
    environment['DMS_PPTX_SOURCE'] = source_path
    environment['DMS_PPTX_DESTINATION'] = temporary_path

    try:
        result = subprocess.run(
            ['powershell.exe', '-NoLogo', '-NoProfile', '-NonInteractive', '-EncodedCommand', encoded_script],
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
            env=environment,
        )
        if result.returncode != 0:
            error_message = result.stderr.strip() or result.stdout.strip() or 'PowerPoint conversion failed.'
            raise RuntimeError(error_message)
        if not os.path.isfile(temporary_path) or os.path.getsize(temporary_path) == 0:
            raise RuntimeError('PowerPoint did not create a PDF preview.')

        os.replace(temporary_path, destination_path)
    finally:
        if os.path.exists(temporary_path):
            os.remove(temporary_path)


def _pptx_to_text_pdf(src_path, dest_path):
    """Create a basic text preview when a native presentation renderer is unavailable."""
    from pptx import Presentation
    prs = Presentation(src_path)
    pdf_doc = SimpleDocTemplate(dest_path, pagesize=A4,
                                 leftMargin=inch, rightMargin=inch,
                                 topMargin=inch, bottomMargin=inch)
    styles = getSampleStyleSheet()
    story = []
    for i, slide in enumerate(prs.slides, 1):
        story.append(Paragraph(f"Slide {i}", styles['Heading2']))
        story.append(Spacer(1, 6))
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                try:
                    story.append(Paragraph(shape.text.strip(), styles['Normal']))
                except Exception:
                    safe = shape.text.strip().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    story.append(Paragraph(safe, styles['Normal']))
                story.append(Spacer(1, 4))
        story.append(Spacer(1, 16))
    if not story:
        story.append(Paragraph("(Empty presentation)", styles['Normal']))
    pdf_doc.build(story)
    return dest_path, None
